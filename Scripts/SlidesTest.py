
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a new presentation object
prs = Presentation()
slide_titles = []
# Define slide titles and contents
slide_titles = [
    "The Future of Human Labor", "Current Workforce Trends", "Technological Impact",
    "Human Labor Limitations", "Automation and Jobs", "Adaptation Strategies",
    "Skills for the Future", "Policy and Regulation", "Real-World Examples", "Concluding Thoughts"
]

slide_contents = [
    "Exploring how the labor market will evolve in the coming years.",
    "Analysis of current trends shaping the global workforce.",
    "The role of AI, robotics, and other technologies in work.",
    "Understanding areas where human labor is less effective.",
    "How automation is transforming job markets and skill requirements.",
    "Strategies to adapt to labor market changes and technological advancements.",
    "Identifying key skills that will be crucial in future job markets.",
    "Policy recommendations to support workers and ensure fair labor practices.",
    "Case studies of how different industries are adapting to these changes.",
    "Summarizing the main points and envisioning the future of work."
]

# Loop to create each slide
for title, content in zip(slide_titles, slide_contents):
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    # Set title and content
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    content_placeholder = slide.placeholders[1]
    content_placeholder.text = content

    # Formatting title
    for paragraph in title_placeholder.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(24)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 0, 0)  # Black color

    # Formatting content
    for paragraph in content_placeholder.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = Pt(18)

# Save the presentation
prs.save('FutureOfHumanLabor.pptx')